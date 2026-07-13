"""Create a 10-slide PowerPoint summary of the neuromotor project."""

from pathlib import Path
import sys

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from src.metrics import force_summary, neural_summary
from src.muscle import spikes_to_force
from src.renshaw import simulate_renshaw_circuit


OUT_DIR = ROOT / "outputs" / "presentation"
ASSET_DIR = OUT_DIR / "assets"
OUT_FILE = OUT_DIR / "simulacion_circuito_neuromotor_izhikevich.pptx"

NAVY = RGBColor(12, 27, 52)
BLUE = RGBColor(29, 111, 164)
CYAN = RGBColor(32, 190, 207)
CORAL = RGBColor(231, 92, 83)
ORANGE = RGBColor(242, 160, 70)
WHITE = RGBColor(247, 250, 252)
LIGHT = RGBColor(224, 234, 242)
MUTED = RGBColor(151, 170, 190)
DARK = RGBColor(26, 43, 66)
GREEN = RGBColor(80, 190, 145)


def run_experiments():
    base = dict(
        n_motor=20, n_renshaw=5, amplitude=12.0, input_type="motor_plan",
        parameter_noise=0.03, input_noise=0.5, p_mn_to_r=0.5,
        w_mn_to_r=1.5, p_r_to_mn=0.6, synaptic_tau=10.0,
        total_time=1000.0, dt=0.5, seed=42,
    )
    scenarios = {
        "Sin": dict(w_r_to_mn=0.0, recurrent_inhibition=False),
        "Débil": dict(w_r_to_mn=0.5, recurrent_inhibition=True),
        "Moderada": dict(w_r_to_mn=1.5, recurrent_inhibition=True),
        "Fuerte": dict(w_r_to_mn=4.0, recurrent_inhibition=True),
    }
    results, metrics, forces = {}, {}, {}
    for name, changes in scenarios.items():
        result = simulate_renshaw_circuit(**base, **changes)
        force = spikes_to_force(result["spikes_MN"], result["time"])
        results[name] = result
        forces[name] = force
        metrics[name] = {
            **neural_summary(result["spikes_MN"], result["time"], result["I_R_to_MN"]),
            **force_summary(force, result["time"]),
        }
    return results, metrics, forces


def set_mpl_style():
    plt.rcParams.update({
        "figure.facecolor": "#0c1b34", "axes.facecolor": "#0c1b34",
        "axes.edgecolor": "#97aabe", "axes.labelcolor": "#e0eaf2",
        "xtick.color": "#97aabe", "ytick.color": "#97aabe",
        "text.color": "#f7fafc", "font.family": "DejaVu Sans",
        "grid.color": "#284363", "grid.alpha": 0.45,
    })


def make_assets(results, metrics, forces):
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    set_mpl_style()

    moderate = results["Moderada"]
    fig, axes = plt.subplots(2, 1, figsize=(10, 5.3), sharex=True)
    axes[0].plot(moderate["time"], moderate["V_MN"][0], color="#20becf", lw=1)
    axes[0].set_ylabel("v (mV)"); axes[0].set_title("Potencial de membrana — MN representativa", loc="left")
    for row in range(moderate["spikes_MN"].shape[0]):
        times = moderate["time"][moderate["spikes_MN"][row]]
        axes[1].vlines(times, row + 0.6, row + 1.4, color="#f7fafc", lw=0.65)
    axes[1].set(xlabel="Tiempo (ms)", ylabel="Motoneurona", ylim=(0.5, 20.5))
    for ax in axes: ax.grid(True)
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "voltage_raster.png", dpi=180, transparent=True)
    plt.close(fig)

    fig, axes = plt.subplots(1, 2, figsize=(9.6, 3.8))
    for ax, matrix, title in [
        (axes[0], moderate["W_MN_to_R"], "W MN → R (+)"),
        (axes[1], moderate["W_R_to_MN"], "W R → MN (magnitud +)"),
    ]:
        image = ax.imshow(matrix, aspect="auto", origin="lower", cmap="viridis")
        ax.set(title=title, xlabel="Destino", ylabel="Origen")
        fig.colorbar(image, ax=ax, fraction=.05)
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "weight_matrices.png", dpi=180, transparent=True)
    plt.close(fig)

    names = list(metrics)
    spikes = [metrics[n]["spikes_totales"] for n in names]
    rates = [metrics[n]["firing_rate_medio_Hz"] for n in names]
    inhibition = [metrics[n]["corriente_inhibitoria_media"] for n in names]
    fig, axes = plt.subplots(1, 3, figsize=(12.5, 3.6))
    colors = ["#97aabe", "#20becf", "#f2a046", "#e75c53"]
    for ax, values, title, ylabel in [
        (axes[0], spikes, "Spikes MN", "Total"),
        (axes[1], rates, "Firing rate medio", "Hz"),
        (axes[2], inhibition, "Corriente inhibitoria", "u.a."),
    ]:
        ax.bar(names, values, color=colors, width=.68)
        ax.set(title=title, ylabel=ylabel); ax.grid(axis="y")
        ax.tick_params(axis="x", rotation=20)
        for index, value in enumerate(values):
            label = f"{value:.2f}" if isinstance(value, float) else str(value)
            ax.text(index, value, label, ha="center", va="bottom", fontsize=9)
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "neural_results.png", dpi=180, transparent=True)
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(10.5, 4.4))
    palette = ["#97aabe", "#20becf", "#f2a046", "#e75c53"]
    for (name, force), color in zip(forces.items(), palette):
        ax.plot(results[name]["time"], force, label=name, color=color, lw=2)
    ax.set(title="Suma de twitches del pool motor", xlabel="Tiempo (ms)", ylabel="Fuerza (u.a.)")
    ax.legend(frameon=False, ncol=4); ax.grid(True)
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "force_comparison.png", dpi=180, transparent=True)
    plt.close(fig)


def add_full_bg(slide, color=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_text(slide, text, x, y, w, h, size=20, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font="Aptos", margin=0.05, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear(); box.text_frame.margin_left = Inches(margin); box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin); box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = valign
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = text; paragraph.alignment = align
    paragraph.font.name = font; paragraph.font.size = Pt(size); paragraph.font.bold = bold; paragraph.font.color.rgb = color
    return box


def add_rich_lines(slide, lines, x, y, w, h, size=17, color=LIGHT, bullet=True, spacing=9):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame; frame.clear(); frame.word_wrap = True
    for index, line in enumerate(lines):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = line; p.font.name = "Aptos"; p.font.size = Pt(size); p.font.color.rgb = color
        p.space_after = Pt(spacing); p.level = 0
        if bullet: p.text = "•  " + line
    return box


def add_title(slide, number, title, subtitle=None):
    add_text(slide, f"{number:02d}", .5, .35, .7, .45, 13, CYAN, True)
    add_text(slide, title, 1.15, .27, 11.55, .55, 25, WHITE, True)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.5), Inches(.91), Inches(12.3), Inches(.025)).fill.solid()
    line = slide.shapes[-1]; line.fill.fore_color.rgb = RGBColor(45, 70, 98); line.line.fill.background()
    if subtitle: add_text(slide, subtitle, 1.15, .86, 11.3, .38, 10, MUTED)


def add_footer(slide, text="Neural Network Exploration · Python/Jupyter"):
    add_text(slide, text, .55, 7.12, 8, .22, 8, MUTED)


def add_pill(slide, text, x, y, w, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(.43))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    add_text(slide, text, x, y+.01, w, .38, 10, WHITE, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_metric_card(slide, value, label, x, y, w=2.3, accent=CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.05))
    shape.fill.solid(); shape.fill.fore_color.rgb = DARK; shape.line.color.rgb = accent
    add_text(slide, value, x+.12, y+.12, w-.24, .42, 22, accent, True, PP_ALIGN.CENTER)
    add_text(slide, label, x+.1, y+.59, w-.2, .3, 9, LIGHT, False, PP_ALIGN.CENTER)


def add_node(slide, label, x, y, w, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(.72))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    add_text(slide, label, x, y+.06, w, .55, 13, WHITE, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return shape


def connect(slide, x1, y1, x2, y2, color=CYAN, width=2.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color; line.line.width = Pt(width); line.line.end_arrowhead = True
    return line


def add_image(slide, path, x, y, w, h=None):
    kwargs = {"left": Inches(x), "top": Inches(y), "width": Inches(w)}
    if h is not None: kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), **kwargs)


def build_presentation(results, metrics, forces):
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 — Cover
    slide = prs.slides.add_slide(blank); add_full_bg(slide)
    for x, y, r, color in [(10.3, 1.0, 1.3, CYAN), (11.5, 2.0, .8, CORAL), (9.6, 3.0, .65, ORANGE), (11.0, 4.3, 1.0, BLUE)]:
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(r), Inches(r))
        node.fill.solid(); node.fill.fore_color.rgb = color; node.fill.transparency = 15; node.line.fill.background()
    connect(slide, 10.8, 1.6, 11.8, 2.25, LIGHT, 1.3); connect(slide, 11.7, 2.8, 11.3, 4.35, LIGHT, 1.3)
    add_pill(slide, "SIMULACIÓN NEUROMOTORA", .7, .65, 2.7, BLUE)
    add_text(slide, "Simulación progresiva de un\ncircuito neuromotor", .7, 1.55, 8.4, 1.5, 34, WHITE, True)
    add_text(slide, "mediante neuronas de Izhikevich", .73, 3.1, 7.7, .55, 25, CYAN, True)
    add_text(slide, "De una neurona individual a una salida muscular simplificada", .73, 4.05, 7.5, .8, 17, LIGHT)
    add_text(slide, "Proyecto Python · Jupyter · Julio 2026", .73, 6.45, 6, .35, 11, MUTED)

    # 2 — Objective
    slide = prs.slides.add_slide(blank); add_full_bg(slide); add_title(slide, 2, "Problema, objetivo y alcance")
    add_text(slide, "¿Cómo modifica la inhibición recurrente la actividad de un pool motor y su salida muscular estimada?", .7, 1.35, 11.8, 1.0, 25, WHITE, True, PP_ALIGN.CENTER)
    add_metric_card(slide, "01", "Actividad spiking", .9, 3.0, 2.7, CYAN)
    add_metric_card(slide, "02", "Feedback Renshaw", 4.05, 3.0, 2.7, CORAL)
    add_metric_card(slide, "03", "Fuerza simplificada", 7.2, 3.0, 2.7, ORANGE)
    add_metric_card(slide, "04", "Interpretación prudente", 10.35, 3.0, 2.2, GREEN)
    add_rich_lines(slide, ["Modelo conceptual y reproducible", "Simulaciones pequeñas y rápidas", "No es una reconstrucción biológica ni biomecánica"], 2.15, 4.7, 9.2, 1.45, 15, LIGHT, True)
    add_footer(slide)

    # 3 — Theory
    slide = prs.slides.add_slide(blank); add_full_bg(slide); add_title(slide, 3, "Marco teórico: modelo de Izhikevich")
    eq = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.7), Inches(1.35), Inches(6.05), Inches(2.0))
    eq.fill.solid(); eq.fill.fore_color.rgb = DARK; eq.line.color.rgb = BLUE
    add_text(slide, "dv/dt = 0.04v² + 5v + 140 − u + I", 1.0, 1.72, 5.45, .55, 23, CYAN, True, PP_ALIGN.CENTER, "Cambria Math")
    add_text(slide, "du/dt = a(bv − u)", 1.0, 2.45, 5.45, .5, 23, WHITE, True, PP_ALIGN.CENTER, "Cambria Math")
    reset = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.7), Inches(3.75), Inches(6.05), Inches(1.45))
    reset.fill.solid(); reset.fill.fore_color.rgb = RGBColor(51, 34, 48); reset.line.color.rgb = CORAL
    add_text(slide, "Si v ≥ 30 mV → v = c ; u = u + d", 1.0, 4.17, 5.45, .5, 19, CORAL, True, PP_ALIGN.CENTER, "Cambria Math")
    add_rich_lines(slide, ["v: potencial de membrana", "u: recuperación", "I: corriente total", "a, b: dinámica de recuperación", "c, d: reset tras el spike"], 7.35, 1.35, 4.9, 3.2, 17, LIGHT)
    add_text(slide, "Ventaja", 7.45, 5.05, 1.2, .35, 12, CYAN, True)
    add_text(slide, "Patrones de disparo ricos con pocas variables y bajo costo computacional.", 7.45, 5.43, 4.8, .8, 16, WHITE)
    add_footer(slide)

    # 4 — Progression
    slide = prs.slides.add_slide(blank); add_full_bg(slide); add_title(slide, 4, "Construcción progresiva del proyecto")
    labels = ["Neurona\nindividual", "Pool\nheterogéneo", "Matriz W", "MN +\nRenshaw", "Inhibición\nrecurrente", "Fuerza\nmuscular"]
    colors = [BLUE, BLUE, CYAN, ORANGE, CORAL, GREEN]
    xs = [.55, 2.6, 4.65, 6.7, 8.75, 10.8]
    for i, (label, x, color) in enumerate(zip(labels, xs, colors)):
        add_node(slide, label, x, 2.35, 1.55, color)
        if i < len(labels)-1: connect(slide, x+1.56, 2.71, xs[i+1]-.05, 2.71, MUTED, 1.8)
    add_text(slide, "01", .95, 1.75, .6, .35, 12, MUTED, True, PP_ALIGN.CENTER)
    for i in range(1, 6): add_text(slide, f"{i+1:02d}", xs[i]+.48, 1.75, .6, .35, 12, MUTED, True, PP_ALIGN.CENTER)
    add_rich_lines(slide, ["Notebooks 01–06: recorrido explicativo e interactivo", "src/: dinámica, inputs, conectividad, circuito, métricas, músculo y visualización", "Semillas fijas para comparar escenarios con la misma realización aleatoria"], 1.0, 4.35, 11.3, 1.7, 16, LIGHT)
    add_footer(slide)

    # 5 — Network
    slide = prs.slides.add_slide(blank); add_full_bg(slide); add_title(slide, 5, "Red conectada y dinámica sináptica")
    add_image(slide, ASSET_DIR / "weight_matrices.png", .45, 1.25, 7.3)
    add_text(slide, "W[i, j]", 8.15, 1.45, 2.1, .55, 24, CYAN, True)
    add_text(slide, "efecto de la neurona i sobre la neurona j", 8.15, 2.0, 4.4, .65, 15, LIGHT)
    add_rich_lines(slide, ["Conectividad Bernoulli aleatoria", "Variable sináptica por neurona", "Decaimiento: exp(−dt/τ)", "Ruido individual de input", "Heterogeneidad en a, b, c y d"], 8.0, 2.9, 4.5, 2.7, 15, LIGHT)
    add_pill(slide, "Filas = origen · Columnas = destino", 8.1, 5.85, 4.0, BLUE)
    add_footer(slide)

    # 6 — Circuit
    slide = prs.slides.add_slide(blank); add_full_bg(slide); add_title(slide, 6, "Circuito de inhibición recurrente MN–Renshaw")
    motor = add_node(slide, "Pool de motoneuronas\nRegular Spiking", 1.05, 2.25, 3.0, BLUE)
    renshaw = add_node(slide, "Células de Renshaw\nFast Spiking*", 8.95, 2.25, 3.0, ORANGE)
    add_node(slide, "Comando motor", 1.55, 4.3, 2.0, DARK)
    connect(slide, 2.55, 4.3, 2.55, 3.0, CYAN, 3)
    connect(slide, 4.05, 2.48, 8.95, 2.48, CYAN, 3)
    connect(slide, 8.95, 2.86, 4.05, 2.86, CORAL, 3)
    add_text(slide, "Excitación MN → R", 5.15, 1.78, 2.8, .35, 13, CYAN, True, PP_ALIGN.CENTER)
    add_text(slide, "Inhibición R → MN", 5.15, 3.15, 2.8, .35, 13, CORAL, True, PP_ALIGN.CENTER)
    add_text(slide, "I_MN = I_motor + I_ruido − I_Renshaw", .8, 5.45, 5.9, .45, 17, WHITE, True, PP_ALIGN.CENTER, "Cambria Math")
    add_text(slide, "I_R = I_MN→R + I_ruido,R", 6.75, 5.45, 5.4, .45, 17, WHITE, True, PP_ALIGN.CENTER, "Cambria Math")
    add_text(slide, "* Aproximación funcional; no calibración biológica específica.", 7.65, 6.45, 4.6, .3, 9, MUTED, False, PP_ALIGN.RIGHT)
    add_footer(slide)

    # 7 — Experiment
    slide = prs.slides.add_slide(blank); add_full_bg(slide); add_title(slide, 7, "Diseño experimental controlado")
    scenario_names = ["Sin inhibición", "Débil", "Moderada", "Fuerte"]
    scenario_weights = ["0,0", "0,5", "1,5", "4,0"]
    card_colors = [MUTED, CYAN, ORANGE, CORAL]
    for i, (name, weight, color) in enumerate(zip(scenario_names, scenario_weights, card_colors)):
        x = .7 + i*3.13
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.55), Inches(2.6), Inches(1.55))
        shape.fill.solid(); shape.fill.fore_color.rgb = DARK; shape.line.color.rgb = color
        add_text(slide, name, x+.1, 1.8, 2.4, .4, 15, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, f"peso R→MN = {weight}", x+.1, 2.35, 2.4, .35, 11, color, True, PP_ALIGN.CENTER)
    add_text(slide, "Todo lo demás permanece constante", .9, 3.85, 11.5, .55, 21, CYAN, True, PP_ALIGN.CENTER)
    constants = [("20", "motoneuronas"), ("5", "Renshaw"), ("1000 ms", "duración"), ("0,5 ms", "dt"), ("42", "semilla")]
    for i, (value, label) in enumerate(constants):
        add_metric_card(slide, value, label, 1.05+i*2.37, 4.75, 2.0, [BLUE, ORANGE, GREEN, CYAN, CORAL][i])
    add_text(slide, "Mismo input · mismo ruido · misma conectividad aleatoria", 2.2, 6.25, 8.9, .4, 13, LIGHT, False, PP_ALIGN.CENTER)
    add_footer(slide)

    # 8 — Neural results
    slide = prs.slides.add_slide(blank); add_full_bg(slide); add_title(slide, 8, "Resultados: regulación de la actividad motoneuronal")
    add_image(slide, ASSET_DIR / "neural_results.png", .35, 1.2, 8.6)
    add_metric_card(slide, "−7,1%", "spikes: 311 → 289", 9.35, 1.55, 2.8, CYAN)
    add_metric_card(slide, "−7,1%", "firing: 15,55 → 14,45 Hz", 9.35, 2.95, 2.8, CORAL)
    add_metric_card(slide, "20/20", "MN activas en todos", 9.35, 4.35, 2.8, GREEN)
    add_text(slide, "La inhibición aumentó y la actividad total bajó modestamente. El pico poblacional no fue monotónico: no se fuerza una conclusión de “mayor estabilidad”.", 1.0, 6.05, 11.4, .75, 14, LIGHT, False, PP_ALIGN.CENTER)
    add_footer(slide)

    # 9 — Muscle
    slide = prs.slides.add_slide(blank); add_full_bg(slide); add_title(slide, 9, "De spikes a fuerza muscular simplificada")
    add_image(slide, ASSET_DIR / "force_comparison.png", .45, 1.18, 8.0)
    add_text(slide, "h(t) = A(e⁻ᵗ/τᵈᵉᶜᵃʸ − e⁻ᵗ/τʳⁱˢᵉ)", 8.65, 1.55, 4.1, .55, 19, CYAN, True, PP_ALIGN.CENTER, "Cambria Math")
    add_rich_lines(slide, ["Cada spike genera un twitch", "La fuerza suma todos los twitches", "Convolución causal del pool completo"], 8.75, 2.45, 3.8, 1.65, 14, LIGHT)
    add_metric_card(slide, "23,22", "media sin inhibición", 8.7, 4.45, 1.75, CYAN)
    add_metric_card(slide, "21,58", "media inhibición fuerte", 10.65, 4.45, 1.75, CORAL)
    add_text(slide, "Proxy visual de activación: no incluye biomecánica, fatiga ni relación longitud–tensión.", 8.75, 5.9, 3.7, .75, 11, MUTED, False, PP_ALIGN.CENTER)
    add_footer(slide)

    # 10 — Conclusions
    slide = prs.slides.add_slide(blank); add_full_bg(slide); add_title(slide, 10, "Conclusiones, limitaciones y próximos pasos")
    columns = [
        ("CONCLUSIONES", CYAN, ["Izhikevich ofrece dinámica rica con bajo costo", "W explicita dirección y magnitud sináptica", "La inhibición redujo actividad y fuerza media en esta ejecución"]),
        ("LIMITACIONES", CORAL, ["Parámetros no calibrados", "Conexiones e input artificiales", "Sin feedback sensorial ni antagonista", "Músculo no biomecánico"]),
        ("TRABAJO FUTURO", GREEN, ["Calibrar con bibliografía y datos", "Agonista–antagonista e inhibición recíproca", "Huso y órgano tendinoso de Golgi", "Feedback y biomecánica"]),
    ]
    for i, (heading, color, items) in enumerate(columns):
        x = .55 + i*4.25
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.35), Inches(3.85), Inches(4.9))
        shape.fill.solid(); shape.fill.fore_color.rgb = DARK; shape.line.color.rgb = color
        add_text(slide, heading, x+.15, 1.65, 3.55, .45, 14, color, True, PP_ALIGN.CENTER)
        add_rich_lines(slide, items, x+.28, 2.35, 3.3, 3.3, 14, LIGHT, True, 14)
    add_text(slide, "La simulación es una plataforma exploratoria: sus resultados describen este modelo y esta parametrización.", 1.15, 6.55, 11.0, .45, 14, WHITE, True, PP_ALIGN.CENTER)
    add_footer(slide, "Neural Network Exploration · Fin")

    prs.core_properties.title = "Simulación progresiva de un circuito neuromotor mediante neuronas de Izhikevich"
    prs.core_properties.subject = "Modelo neuronal, circuito Renshaw y fuerza muscular simplificada"
    prs.core_properties.author = "Neural Network Exploration"
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)
    return OUT_FILE


if __name__ == "__main__":
    results, metrics, forces = run_experiments()
    make_assets(results, metrics, forces)
    path = build_presentation(results, metrics, forces)
    print(path)
